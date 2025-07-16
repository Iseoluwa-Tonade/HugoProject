import streamlit as st
from openai import OpenAI # <-- CHANGE: Import OpenAI
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import tempfile

# --- Core Functions for Text Extraction (No changes here) ---

def extract_text_from_pdf(filepath):
    text = ""
    with open(filepath, 'rb') as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_pptx(filepath):
    text = ""
    prs = Presentation(filepath)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(filepath):
    text = ""
    doc = Document(filepath)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def process_uploaded_files(uploaded_files):
    combined_text = ""
    for uploaded_file in uploaded_files:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_filepath = tmp_file.name

        st.info(f"Processing {uploaded_file.name}...")
        try:
            if tmp_filepath.endswith('.pdf'):
                combined_text += extract_text_from_pdf(tmp_filepath) + "\n\n"
            elif tmp_filepath.endswith('.pptx'):
                combined_text += extract_text_from_pptx(tmp_filepath) + "\n\n"
            elif tmp_filepath.endswith('.docx'):
                combined_text += extract_text_from_docx(tmp_filepath) + "\n\n"
        finally:
            os.remove(tmp_filepath)

    return combined_text

# --- Streamlit App UI and Logic ---

st.set_page_config(page_title="Query Your Docs", page_icon="📄")
st.title("📄 Query Your Documents")

# Sidebar for API Key and file uploads
with st.sidebar:
    st.header("Configuration")
    try:
        # <-- CHANGE: Use OpenAI API Key from secrets
        api_key = st.secrets["OPENAI_API_KEY"]
        st.success("API key loaded from secrets!")
    except:
        # <-- CHANGE: Update text for OpenAI
        api_key = st.text_input("Enter your OpenAI API Key:", type="password")

    uploaded_files = st.file_uploader(
        "Upload your documents (PDF, DOCX, PPTX)",
        type=['pdf', 'docx', 'pptx'],
        accept_multiple_files=True
    )

# Main app logic
if not api_key:
    st.warning("Please enter your OpenAI API Key in the sidebar to proceed.")
    st.stop()

# <-- CHANGE: Initialize OpenAI client
try:
    client = OpenAI(api_key=api_key)
except Exception as e:
    st.error(f"Failed to initialize OpenAI client: {e}")
    st.stop()

# Process files only once and store in session state
if "processed_text" not in st.session_state:
    st.session_state.processed_text = None

if uploaded_files:
    if st.button("Process Documents"):
        with st.spinner("Processing documents... this may take a moment."):
            st.session_state.processed_text = process_uploaded_files(uploaded_files)
        st.success("Documents processed successfully! You can now ask questions.")

# Initialize chat history
if "messages" not in st.session_state:
    st.session_state.messages = []

# Display chat messages from history
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Accept user input
if prompt := st.chat_input("Ask a question about your documents"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if st.session_state.processed_text is None:
        st.warning("Please upload and process documents before asking a question.")
    else:
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                # <-- CHANGE: Structure the prompt and call the OpenAI API
                try:
                    # Construct the messages list for OpenAI
                    messages = [
                        {"role": "system", "content": "You are a helpful assistant that answers questions based on the provided document text."},
                        {"role": "user", "content": f"""
                            Based on the following text extracted from the user's documents, please answer their question.
                            If the answer is not found in the text, state that clearly.

                            ---
                            DOCUMENT TEXT:
                            {st.session_state.processed_text}
                            ---

                            USER'S QUESTION:
                            {prompt}
                        """}
                    ]

                    # Call the OpenAI API
                    response = client.chat.completions.create(
                        model="gpt-4o",  # Or "gpt-3.5-turbo"
                        messages=messages
                    )
                    response_text = response.choices[0].message.content
                except Exception as e:
                    response_text = f"An error occurred with the OpenAI API: {e}"

                st.markdown(response_text)
        st.session_state.messages.append({"role": "assistant", "content": response_text})
